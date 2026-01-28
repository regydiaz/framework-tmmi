"""
Módulo de exportação do Framework TMMi
Gera relatórios em PDF e apresentações em PowerPoint
"""

import pandas as pd
from datetime import datetime
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io


class TMMiExporter:
    """Classe para exportar dados do TMMi para PDF e PowerPoint"""
    
    def __init__(self, data_dict):
        """
        Inicializa o exportador
        
        Args:
            data_dict: Dicionário com os dataframes carregados
        """
        self.data = data_dict
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        """Configura estilos customizados"""
        # Título principal
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Subtítulo
        self.styles.add(ParagraphStyle(
            name='CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#333333'),
            spaceAfter=12,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        ))
        
        # Texto normal
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=10,
            spaceAfter=6,
            fontName='Helvetica'
        ))
    
    def export_to_pdf(self, output_path='/mnt/user-data/outputs/Framework_TMMi_Relatorio.pdf'):
        """
        Exporta relatório completo em PDF
        
        Args:
            output_path: Caminho do arquivo de saída
        
        Returns:
            str: Caminho do arquivo gerado
        """
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        story = []
        
        # Título
        title = Paragraph("Framework TMMi - TAG IMF", self.styles['CustomTitle'])
        story.append(title)
        
        subtitle = Paragraph(
            f"Relatório Executivo - {datetime.now().strftime('%d/%m/%Y')}",
            self.styles['CustomBody']
        )
        story.append(subtitle)
        story.append(Spacer(1, 0.3*inch))
        
        # ===== VISÃO INSTITUCIONAL =====
        story.append(Paragraph("Visão Institucional", self.styles['CustomHeading']))
        
        df_inst = self.data['institucional']
        
        # Estatísticas gerais
        total_areas = len(df_inst)
        adotado = len(df_inst[df_inst['Status Institucional'] == 'Adotado'])
        desenvolvendo = len(df_inst[df_inst['Status Institucional'] == 'Desenvolvendo'])
        
        stats_text = f"""
        <b>Áreas de Processo Mapeadas:</b> {total_areas}<br/>
        <b>Adotado:</b> {adotado} ({(adotado/total_areas*100):.0f}%)<br/>
        <b>Desenvolvendo:</b> {desenvolvendo} ({(desenvolvendo/total_areas*100):.0f}%)
        """
        story.append(Paragraph(stats_text, self.styles['CustomBody']))
        story.append(Spacer(1, 0.2*inch))
        
        # Tabela de status por nível
        data_table = [['Nível', 'Área de Processo', 'Status', 'Observação']]
        
        for idx, row in df_inst.iterrows():
            nivel = str(row.get('Nível TMMi', ''))
            area = str(row.get('Área de Processo', ''))
            status = str(row.get('Status Institucional', ''))
            obs = str(row.get('Observação', ''))[:50] + '...' if len(str(row.get('Observação', ''))) > 50 else str(row.get('Observação', ''))
            
            if area and area != 'nan':
                data_table.append([nivel, area, status, obs])
        
        if len(data_table) > 1:
            table = Table(data_table, colWidths=[0.8*inch, 2*inch, 1.2*inch, 2*inch])
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey])
            ]))
            story.append(table)
        
        story.append(PageBreak())
        
        # ===== ROADMAP TRIMESTRAL =====
        story.append(Paragraph("Roadmap Trimestral", self.styles['CustomHeading']))
        
        df_roadmap = self.data['roadmap']
        
        if not df_roadmap.empty:
            roadmap_table = [['Trimestre', 'Fase', 'Entrega', 'Status']]
            
            for idx, row in df_roadmap.iterrows():
                tri = str(row.get('Trimestre', ''))[:15]
                fase = str(row.get('Fase', ''))[:20]
                entrega = str(row.get('Entrega', ''))[:50] + '...' if len(str(row.get('Entrega', ''))) > 50 else str(row.get('Entrega', ''))
                status = str(row.get('Status', ''))
                
                if entrega and entrega != 'nan':
                    roadmap_table.append([tri, fase, entrega, status])
            
            if len(roadmap_table) > 1:
                table = Table(roadmap_table, colWidths=[1*inch, 1.3*inch, 3*inch, 1*inch])
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 1), (-1, -1), 8),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey])
                ]))
                story.append(table)
        
        story.append(PageBreak())
        
        # ===== MAPA DO TMMi =====
        story.append(Paragraph("Mapa do TMMi", self.styles['CustomHeading']))
        
        df_mapa = self.data['mapa']
        
        if not df_mapa.empty and 'Nível' in df_mapa.columns:
            niveis = df_mapa['Nível'].dropna().unique()
            
            for nivel in sorted(niveis):
                df_nivel = df_mapa[df_mapa['Nível'] == nivel]
                
                story.append(Paragraph(f"<b>Nível {nivel}</b>", self.styles['CustomBody']))
                story.append(Spacer(1, 0.1*inch))
                
                for idx, row in df_nivel.iterrows():
                    area = row.get('Área de Processo', 'N/A')
                    descricao = str(row.get('Descrição', ''))[:200] + '...' if len(str(row.get('Descrição', ''))) > 200 else str(row.get('Descrição', ''))
                    
                    if area != 'N/A' and str(area) != 'nan':
                        story.append(Paragraph(f"• <b>{area}</b>", self.styles['CustomBody']))
                        if descricao and descricao != 'nan':
                            story.append(Paragraph(f"  {descricao}", self.styles['CustomBody']))
                        story.append(Spacer(1, 0.1*inch))
                
                story.append(Spacer(1, 0.2*inch))
        
        # Gerar PDF
        doc.build(story)
        return output_path
    
    def export_to_powerpoint(self, output_path='/mnt/user-data/outputs/Framework_TMMi_Apresentacao.pptx'):
        """
        Exporta apresentação executiva em PowerPoint
        
        Args:
            output_path: Caminho do arquivo de saída
        
        Returns:
            str: Caminho do arquivo gerado
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # ===== SLIDE 1: TÍTULO =====
        slide_title = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Título
        title_box = slide_title.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Framework TMMi"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 119, 180)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtítulo
        subtitle_box = slide_title.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"TAG IMF - {datetime.now().strftime('%B %Y')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # ===== SLIDE 2: VISÃO GERAL =====
        slide_overview = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        title = slide_overview.shapes.title
        title.text = "Visão Geral do TMMi"
        
        df_inst = self.data['institucional']
        total_areas = len(df_inst)
        adotado = len(df_inst[df_inst['Status Institucional'] == 'Adotado'])
        desenvolvendo = len(df_inst[df_inst['Status Institucional'] == 'Desenvolvendo'])
        em_adocao = len(df_inst[df_inst['Status Institucional'] == 'Em Adoção'])
        
        # Criar caixas de métricas
        metrics = [
            ("Total de Áreas", total_areas, 1, 2),
            ("Adotado", adotado, 3.5, 2),
            ("Desenvolvendo", desenvolvendo, 6, 2),
            ("Em Adoção", em_adocao, 1, 4.5)
        ]
        
        for metric_name, metric_value, left, top in metrics:
            # Box de fundo
            box = slide_overview.shapes.add_shape(
                1,  # Rectangle
                Inches(left), Inches(top),
                Inches(2), Inches(1.2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(227, 242, 253)
            box.line.color.rgb = RGBColor(31, 119, 180)
            
            # Valor
            value_box = slide_overview.shapes.add_textbox(
                Inches(left), Inches(top + 0.1),
                Inches(2), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = str(metric_value)
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(36)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(31, 119, 180)
            value_para.alignment = PP_ALIGN.CENTER
            
            # Nome da métrica
            name_box = slide_overview.shapes.add_textbox(
                Inches(left), Inches(top + 0.7),
                Inches(2), Inches(0.4)
            )
            name_frame = name_box.text_frame
            name_frame.text = metric_name
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(14)
            name_para.alignment = PP_ALIGN.CENTER
        
        # ===== SLIDE 3: STATUS POR NÍVEL =====
        slide_status = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide_status.shapes.title
        title.text = "Status por Nível TMMi"
        
        # Tabela de status
        rows = 1 + len(df_inst[df_inst['Área de Processo'].notna()])
        cols = 3
        
        table = slide_status.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
        
        # Cabeçalho
        table.cell(0, 0).text = "Nível"
        table.cell(0, 1).text = "Área de Processo"
        table.cell(0, 2).text = "Status"
        
        # Formatar cabeçalho
        for i in range(3):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Dados
        row_idx = 1
        for idx, row in df_inst.iterrows():
            if pd.notna(row.get('Área de Processo')) and row_idx < rows:
                table.cell(row_idx, 0).text = str(row.get('Nível TMMi', ''))
                table.cell(row_idx, 1).text = str(row.get('Área de Processo', ''))
                table.cell(row_idx, 2).text = str(row.get('Status Institucional', ''))
                
                # Colorir status
                status = str(row.get('Status Institucional', '')).lower()
                status_cell = table.cell(row_idx, 2)
                status_cell.fill.solid()
                
                if 'adotado' in status:
                    status_cell.fill.fore_color.rgb = RGBColor(212, 237, 218)
                elif 'desenvolvendo' in status:
                    status_cell.fill.fore_color.rgb = RGBColor(255, 243, 205)
                elif 'adoção' in status or 'adocao' in status:
                    status_cell.fill.fore_color.rgb = RGBColor(226, 227, 229)
                
                row_idx += 1
        
        # ===== SLIDE 4: ROADMAP =====
        slide_roadmap = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide_roadmap.shapes.title
        title.text = "Roadmap - Próximas Entregas"
        
        df_roadmap = self.data['roadmap']
        
        if not df_roadmap.empty:
            # Filtrar TRI 1
            df_tri1 = df_roadmap[df_roadmap['Trimestre'].str.contains('TRI 1', na=False)]
            
            if not df_tri1.empty:
                rows_roadmap = min(len(df_tri1) + 1, 8)  # Limitar a 7 entregas
                cols_roadmap = 3
                
                table_roadmap = slide_roadmap.shapes.add_table(
                    rows_roadmap, cols_roadmap,
                    Inches(0.5), Inches(2),
                    Inches(9), Inches(4.5)
                ).table
                
                # Cabeçalho
                table_roadmap.cell(0, 0).text = "Fase"
                table_roadmap.cell(0, 1).text = "Entrega"
                table_roadmap.cell(0, 2).text = "Status"
                
                for i in range(3):
                    cell = table_roadmap.cell(0, i)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(11)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                
                # Dados
                for idx, (_, row) in enumerate(df_tri1.head(7).iterrows(), 1):
                    if idx < rows_roadmap:
                        fase = str(row.get('Fase', ''))[:30]
                        entrega = str(row.get('Entrega', ''))[:60]
                        status = str(row.get('Status', ''))
                        
                        table_roadmap.cell(idx, 0).text = fase
                        table_roadmap.cell(idx, 1).text = entrega
                        table_roadmap.cell(idx, 2).text = status
                        
                        # Ajustar tamanho da fonte
                        for col in range(3):
                            cell = table_roadmap.cell(idx, col)
                            para = cell.text_frame.paragraphs[0]
                            para.font.size = Pt(9)
        
        # ===== SLIDE 5: PRÓXIMOS PASSOS =====
        slide_next = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide_next.shapes.title
        title.text = "Próximos Passos"
        
        content_box = slide_next.placeholders[1]
        tf = content_box.text_frame
        tf.clear()
        
        next_steps = [
            "Consolidar adoção das áreas em desenvolvimento",
            "Avançar para Nível 3 do TMMi",
            "Expandir automação de testes",
            "Fortalecer governança de qualidade",
            "Medir resultados e ajustar estratégia"
        ]
        
        for step in next_steps:
            p = tf.add_paragraph()
            p.text = step
            p.level = 0
            p.font.size = Pt(20)
        
        # Salvar apresentação
        prs.save(output_path)
        return output_path


def export_framework(data_dict, export_pdf=True, export_ppt=True):
    """
    Função helper para exportar o framework
    
    Args:
        data_dict: Dicionário com os dados
        export_pdf: Se True, gera PDF
        export_ppt: Se True, gera PowerPoint
    
    Returns:
        dict: Caminhos dos arquivos gerados
    """
    exporter = TMMiExporter(data_dict)
    results = {}
    
    if export_pdf:
        pdf_path = exporter.export_to_pdf()
        results['pdf'] = pdf_path
    
    if export_ppt:
        ppt_path = exporter.export_to_powerpoint()
        results['ppt'] = ppt_path
    
    return results
